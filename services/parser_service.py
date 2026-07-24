"""
services/parser_service.py — Servicio de extracción de contenido.

Responsabilidad única: convertir cualquier documento soportado en texto
estructurado. Selecciona automáticamente el parser según el tipo de archivo.

Nunca accede a la base de datos.
Nunca utiliza IA.
"""

import csv
import hashlib
import io
import logging
import subprocess
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


# ─── Estructuras de datos ─────────────────────────────────────────

@dataclass
class PageContent:
    """Contenido de una página o sección del documento."""
    numero_pagina: int
    texto: str
    hash_pagina: str = ""


@dataclass
class ParseResult:
    """Resultado de la extracción de contenido de un documento."""
    texto_completo: str
    paginas: list[PageContent] = field(default_factory=list)
    idioma: str | None = None
    num_paginas: int = 0
    parser_utilizado: str = ""


# ─── Helpers ──────────────────────────────────────────────────────

def _calcular_hash(texto: str) -> str:
    """Calcula el hash SHA256 de un texto."""
    return hashlib.sha256(texto.encode("utf-8")).hexdigest()


def _dividir_en_paginas(
    texto: str,
    tam_pagina: int = 2000,
) -> list[PageContent]:
    """Divide texto sin paginación natural en bloques.

    Args:
        texto: Texto completo a dividir.
        tam_pagina: Tamaño máximo de cada bloque en caracteres.

    Returns:
        Lista de PageContent con páginas numeradas.
    """
    if not texto.strip():
        return []

    paginas: list[PageContent] = []
    inicio = 0
    num = 1

    while inicio < len(texto):
        fin = inicio + tam_pagina

        # Intentar cortar en un salto de línea cercano
        if fin < len(texto):
            corte = texto.rfind("\n", inicio, fin)
            if corte > inicio:
                fin = corte + 1

        bloque = texto[inicio:fin].strip()
        if bloque:
            paginas.append(PageContent(
                numero_pagina=num,
                texto=bloque,
                hash_pagina=_calcular_hash(bloque),
            ))
            num += 1

        inicio = fin

    return paginas


# ─── Parser principal ─────────────────────────────────────────────

class ParserService:
    """Extrae texto estructurado de documentos.

    Selecciona automáticamente el parser adecuado según la extensión.
    Cada formato tiene su parser especializado.
    """

    # Mapeo extensión → método parser
    _PARSERS: dict[str, str] = {
        "pdf": "_parsear_pdf",
        "docx": "_parsear_docx",
        "doc": "_parsear_doc",
        "txt": "_parsear_texto_plano",
        "md": "_parsear_texto_plano",
        "csv": "_parsear_csv",
        "xlsx": "_parsear_xlsx",
        "ods": "_parsear_ods",
        "pptx": "_parsear_pptx",
        "odp": "_parsear_odp",
        "html": "_parsear_html",
        "png": "_parsear_imagen",
        "jpg": "_parsear_imagen",
        "jpeg": "_parsear_imagen",
        "webp": "_parsear_imagen",
    }

    @staticmethod
    def extraer_contenido(ruta: Path, extension: str) -> ParseResult:
        """Extrae contenido textual de un documento.

        Args:
            ruta: Ruta al archivo.
            extension: Extensión del archivo (sin punto, en minúsculas).

        Returns:
            ParseResult con el texto extraído y páginas.
        """
        extension = extension.lower().strip(".")

        nombre_metodo = ParserService._PARSERS.get(extension)
        if not nombre_metodo:
            logger.warning("Extensión no soportada para parsing: '%s'", extension)
            return ParseResult(
                texto_completo="",
                parser_utilizado=f"no_soportado:{extension}",
            )

        metodo = getattr(ParserService, nombre_metodo)

        logger.info(
            "Parseando '%s' con %s...", ruta.name, nombre_metodo,
        )

        try:
            resultado = metodo(ruta)
            resultado.num_paginas = len(resultado.paginas)
            logger.info(
                "Parsing completado: %d caracteres, %d páginas (parser: %s)",
                len(resultado.texto_completo), resultado.num_paginas,
                resultado.parser_utilizado,
            )
            return resultado

        except Exception as e:
            logger.error(
                "Error parseando '%s': %s", ruta.name, e, exc_info=True,
            )
            return ParseResult(
                texto_completo="",
                parser_utilizado=f"error:{nombre_metodo}",
            )

    # ─── Parsers específicos ──────────────────────────────────────

    @staticmethod
    def _parsear_pdf(ruta: Path) -> ParseResult:
        """Extrae texto de un PDF. Usa OCR como fallback si no hay texto."""
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber no instalado. No se puede parsear PDF.")
            return ParseResult(texto_completo="", parser_utilizado="pdfplumber:no_instalado")

        paginas: list[PageContent] = []
        textos: list[str] = []
        parser_nombre = "pdfplumber"

        with pdfplumber.open(ruta) as pdf:
            for i, pagina in enumerate(pdf.pages, start=1):
                texto_pagina = pagina.extract_text() or ""

                # Si no hay texto, intentar OCR
                if not texto_pagina.strip():
                    texto_pagina = _ocr_pagina_pdf(pagina, i)
                    if texto_pagina:
                        parser_nombre = "pdfplumber+ocr"

                if texto_pagina.strip():
                    textos.append(texto_pagina)
                    paginas.append(PageContent(
                        numero_pagina=i,
                        texto=texto_pagina.strip(),
                        hash_pagina=_calcular_hash(texto_pagina.strip()),
                    ))

        texto_completo = "\n\n".join(textos)
        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado=parser_nombre,
        )

    @staticmethod
    def _parsear_docx(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo DOCX."""
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx no instalado. No se puede parsear DOCX.")
            return ParseResult(texto_completo="", parser_utilizado="python-docx:no_instalado")

        doc = Document(str(ruta))
        parrafos = [p.text for p in doc.paragraphs if p.text.strip()]
        texto_completo = "\n".join(parrafos)

        paginas = _dividir_en_paginas(texto_completo)

        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado="python-docx",
        )

    @staticmethod
    def _parsear_doc(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo DOC usando antiword."""
        try:
            resultado = subprocess.run(
                ["antiword", str(ruta)],
                capture_output=True, text=True, timeout=30,
            )
            if resultado.returncode == 0 and resultado.stdout.strip():
                texto = resultado.stdout.strip()
                return ParseResult(
                    texto_completo=texto,
                    paginas=_dividir_en_paginas(texto),
                    parser_utilizado="antiword",
                )
        except FileNotFoundError:
            logger.warning("antiword no disponible. No se puede parsear DOC.")
        except subprocess.TimeoutExpired:
            logger.warning("Timeout parseando DOC con antiword.")
        except Exception as e:
            logger.warning("Error con antiword: %s", e)

        return ParseResult(texto_completo="", parser_utilizado="antiword:no_disponible")

    @staticmethod
    def _parsear_texto_plano(ruta: Path) -> ParseResult:
        """Extrae texto de archivos TXT y Markdown."""
        try:
            texto = ruta.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            texto = ruta.read_text(encoding="latin-1")

        return ParseResult(
            texto_completo=texto,
            paginas=_dividir_en_paginas(texto),
            parser_utilizado="texto_plano",
        )

    @staticmethod
    def _parsear_csv(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo CSV."""
        try:
            contenido = ruta.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            contenido = ruta.read_text(encoding="latin-1")

        reader = csv.reader(io.StringIO(contenido))
        filas: list[str] = []

        for fila in reader:
            filas.append(" | ".join(fila))

        texto_completo = "\n".join(filas)

        # Agrupar filas en páginas (~50 filas por página)
        paginas: list[PageContent] = []
        tam_grupo = 50
        for i in range(0, len(filas), tam_grupo):
            grupo = filas[i:i + tam_grupo]
            texto_grupo = "\n".join(grupo)
            paginas.append(PageContent(
                numero_pagina=(i // tam_grupo) + 1,
                texto=texto_grupo,
                hash_pagina=_calcular_hash(texto_grupo),
            ))

        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado="csv_stdlib",
        )

    @staticmethod
    def _parsear_xlsx(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo XLSX. Cada hoja = una página."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("openpyxl no instalado. No se puede parsear XLSX.")
            return ParseResult(texto_completo="", parser_utilizado="openpyxl:no_instalado")

        wb = load_workbook(str(ruta), read_only=True, data_only=True)
        paginas: list[PageContent] = []
        textos: list[str] = []

        for i, nombre_hoja in enumerate(wb.sheetnames, start=1):
            hoja = wb[nombre_hoja]
            filas: list[str] = []

            for fila in hoja.iter_rows(values_only=True):
                valores = [str(v) if v is not None else "" for v in fila]
                if any(v.strip() for v in valores):
                    filas.append(" | ".join(valores))

            texto_hoja = f"[Hoja: {nombre_hoja}]\n" + "\n".join(filas)
            textos.append(texto_hoja)
            paginas.append(PageContent(
                numero_pagina=i,
                texto=texto_hoja,
                hash_pagina=_calcular_hash(texto_hoja),
            ))

        wb.close()
        texto_completo = "\n\n".join(textos)

        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado="openpyxl",
        )

    @staticmethod
    def _parsear_ods(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo ODS. Cada hoja = una página."""
        try:
            from odf.opendocument import load
            from odf.table import Table, TableRow, TableCell
            from odf.text import P
        except ImportError:
            logger.warning("odfpy no instalado. No se puede parsear ODS.")
            return ParseResult(texto_completo="", parser_utilizado="odfpy:no_instalado")

        doc = load(str(ruta))
        paginas: list[PageContent] = []
        textos: list[str] = []

        for i, tabla in enumerate(doc.spreadsheet.getElementsByType(Table), start=1):
            nombre = tabla.getAttribute("name") or f"Hoja {i}"
            filas_texto: list[str] = []

            for fila in tabla.getElementsByType(TableRow):
                celdas = []
                for celda in fila.getElementsByType(TableCell):
                    textos_celda = []
                    for p in celda.getElementsByType(P):
                        # Extraer texto de cada párrafo dentro de la celda
                        texto_p = ""
                        for nodo in p.childNodes:
                            if hasattr(nodo, "data"):
                                texto_p += nodo.data
                            elif hasattr(nodo, "__str__"):
                                texto_p += str(nodo)
                        textos_celda.append(texto_p)
                    celdas.append(" ".join(textos_celda))

                if any(c.strip() for c in celdas):
                    filas_texto.append(" | ".join(celdas))

            texto_hoja = f"[Hoja: {nombre}]\n" + "\n".join(filas_texto)
            textos.append(texto_hoja)
            paginas.append(PageContent(
                numero_pagina=i,
                texto=texto_hoja,
                hash_pagina=_calcular_hash(texto_hoja),
            ))

        texto_completo = "\n\n".join(textos)
        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado="odfpy",
        )

    @staticmethod
    def _parsear_pptx(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo PPTX. Cada slide = una página."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx no instalado. No se puede parsear PPTX.")
            return ParseResult(texto_completo="", parser_utilizado="python-pptx:no_instalado")

        prs = Presentation(str(ruta))
        paginas: list[PageContent] = []
        textos: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            textos_slide: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for parrafo in shape.text_frame.paragraphs:
                        texto_p = parrafo.text.strip()
                        if texto_p:
                            textos_slide.append(texto_p)

            texto_slide = f"[Slide {i}]\n" + "\n".join(textos_slide)
            textos.append(texto_slide)
            paginas.append(PageContent(
                numero_pagina=i,
                texto=texto_slide,
                hash_pagina=_calcular_hash(texto_slide),
            ))

        texto_completo = "\n\n".join(textos)
        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado="python-pptx",
        )

    @staticmethod
    def _parsear_odp(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo ODP. Cada slide = una página."""
        try:
            from odf.opendocument import load
            from odf.draw import Page, Frame, TextBox
            from odf.text import P
        except ImportError:
            logger.warning("odfpy no instalado. No se puede parsear ODP.")
            return ParseResult(texto_completo="", parser_utilizado="odfpy:no_instalado")

        doc = load(str(ruta))
        paginas: list[PageContent] = []
        textos: list[str] = []

        for i, pagina in enumerate(doc.presentation.getElementsByType(Page), start=1):
            textos_pagina: list[str] = []

            for p in pagina.getElementsByType(P):
                texto_p = ""
                for nodo in p.childNodes:
                    if hasattr(nodo, "data"):
                        texto_p += nodo.data
                    elif hasattr(nodo, "__str__"):
                        texto_p += str(nodo)
                if texto_p.strip():
                    textos_pagina.append(texto_p.strip())

            texto_slide = f"[Slide {i}]\n" + "\n".join(textos_pagina)
            textos.append(texto_slide)
            paginas.append(PageContent(
                numero_pagina=i,
                texto=texto_slide,
                hash_pagina=_calcular_hash(texto_slide),
            ))

        texto_completo = "\n\n".join(textos)
        return ParseResult(
            texto_completo=texto_completo,
            paginas=paginas,
            parser_utilizado="odfpy",
        )

    @staticmethod
    def _parsear_html(ruta: Path) -> ParseResult:
        """Extrae texto de un archivo HTML."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            logger.warning("beautifulsoup4 no instalado. No se puede parsear HTML.")
            return ParseResult(texto_completo="", parser_utilizado="bs4:no_instalado")

        try:
            contenido = ruta.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            contenido = ruta.read_text(encoding="latin-1")

        soup = BeautifulSoup(contenido, "html.parser")

        # Eliminar scripts y estilos
        for tag in soup(["script", "style"]):
            tag.decompose()

        texto = soup.get_text(separator="\n")
        # Limpiar líneas vacías excesivas
        lineas = [l.strip() for l in texto.splitlines() if l.strip()]
        texto_limpio = "\n".join(lineas)

        return ParseResult(
            texto_completo=texto_limpio,
            paginas=_dividir_en_paginas(texto_limpio),
            parser_utilizado="beautifulsoup4",
        )

    @staticmethod
    def _parsear_imagen(ruta: Path) -> ParseResult:
        """Extrae texto de una imagen usando OCR (Tesseract)."""
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            logger.warning("Pillow o pytesseract no instalados. No se puede hacer OCR.")
            return ParseResult(texto_completo="", parser_utilizado="ocr:no_instalado")

        try:
            imagen = Image.open(ruta)
            texto = pytesseract.image_to_string(imagen, lang="spa")
            texto = texto.strip()

            paginas = []
            if texto:
                paginas.append(PageContent(
                    numero_pagina=1,
                    texto=texto,
                    hash_pagina=_calcular_hash(texto),
                ))

            return ParseResult(
                texto_completo=texto,
                paginas=paginas,
                parser_utilizado="pytesseract",
            )

        except Exception as e:
            logger.error("Error en OCR para '%s': %s", ruta.name, e)
            return ParseResult(texto_completo="", parser_utilizado="pytesseract:error")


# ─── Helper OCR para PDFs escaneados ──────────────────────────────

def _ocr_pagina_pdf(pagina, num_pagina: int) -> str:
    """Intenta OCR en una página de PDF que no tiene texto extraíble.

    Args:
        pagina: Página de pdfplumber.
        num_pagina: Número de la página (para logging).

    Returns:
        Texto extraído por OCR o string vacío.
    """
    try:
        from PIL import Image
        import pytesseract

        imagen = pagina.to_image(resolution=200)
        # pdfplumber.to_image() retorna un PageImage, necesitamos la imagen PIL
        pil_image = imagen.original
        texto = pytesseract.image_to_string(pil_image, lang="spa")

        if texto.strip():
            logger.debug("OCR exitoso en página %d: %d caracteres.", num_pagina, len(texto))
            return texto.strip()

    except ImportError:
        logger.debug("OCR no disponible (Pillow/pytesseract no instalados).")
    except Exception as e:
        logger.debug("Error en OCR página %d: %s", num_pagina, e)

    return ""
